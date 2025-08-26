import os
import random
import string
from typing import Dict, List, Tuple
import fitz
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont


def random_suffix(n: int = 6) -> str:
    return ''.join(random.choices(string.ascii_lowercase + string.digits, k=n))

def basename(path: str) -> str:
    return os.path.splitext(os.path.basename(path))[0]


def ensure_dir(path: str) -> None:
    os.makedirs(path, exist_ok=True)


def _extract_pdf_text(pdf_path: str) -> str:
    # Using PyMuPDF (fitz) for faster and richer text extraction
    try:
        doc = fitz.open(pdf_path)
        texts: List[str] = []
        for page in doc:
            try:
                txt = page.get_text("text") or ""
                texts.append(txt)
            except Exception:
                continue
        doc.close()
        raw = "\n".join(texts)
        raw = raw.replace('\n\n', '<PARAGRAPH_BREAK>')
        raw = raw.replace('\n', ' ')
        raw = raw.replace('<PARAGRAPH_BREAK>', '\n\n')
        return raw.strip()
    except Exception:
        return ""


def _extract_pdf_images_fitz(pdf_path: str, output_folder: str) -> List[str]:
    saved: List[str] = []
    try:
        doc = fitz.open(pdf_path)
        base = basename(pdf_path)
        for page_index in range(len(doc)):
            page = doc[page_index]
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list, start=1):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n > 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    filename = f"{base}_p{page_index+1}_i{img_index}_{random_suffix()}.png"
                    out_path = os.path.join(output_folder, filename)
                    pix.save(out_path)
                    saved.append(out_path)
                except Exception:
                    continue
        doc.close()
    except Exception:
        return saved
    return saved


def _parse_pdf(pdf_path: str, upload_folder: str) -> Tuple[str, List[str]]:

    text = _extract_pdf_text(pdf_path)
    images: List[str] = []
    images = _extract_pdf_images_fitz(pdf_path, upload_folder)
    return text, images


def _parse_pptx(pptx_path: str, upload_folder: str) -> Tuple[str, List[str]]:

    extracted_text_parts: List[str] = []
    saved_images: List[str] = []

    try:
        prs = Presentation(pptx_path)
        base = basename(pptx_path)
        for s_index, slide in enumerate(prs.slides, start=1):
            # Extract text from shapes
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in paragraph.runs)
                            if line:
                                extracted_text_parts.append(line)
                    # 13 == PICTURE
                    if getattr(shape, "shape_type", None) == 13:
                        image = shape.image
                        image_bytes = image.blob
                        filename = f"{base}_s{s_index}_pic_{random_suffix()}.png"
                        out_path = os.path.join(upload_folder, filename)
                        with open(out_path, "wb") as f:
                            f.write(image_bytes)
                        saved_images.append(out_path)
                except Exception:
                    continue
    except Exception:
        return "", []

    extracted_text = "\n".join(extracted_text_parts).strip()
    return extracted_text, saved_images


def _parse_image(image_path: str) -> Tuple[str, List[str]]:
    return "", [image_path]


def parse_any(file_path: str, upload_folder: str) -> Dict[str, object]:

    ensure_dir(upload_folder)
    ext = os.path.splitext(file_path)[1].lower()

    text = ""
    images: List[str] = []

    if ext == ".pdf":
        text, images = _parse_pdf(file_path, upload_folder)
    elif ext == ".pptx":
        text, images = _parse_pptx(file_path, upload_folder)
    elif ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp"]:
        text, images = _parse_image(file_path)
    else:
        text, images = "", []

    return {
        "text": text or "",
        "images": images or [],
    }

# ---------- Preview helpers ----------

def render_pdf_previews(pdf_path: str, upload_folder: str, pages: int = 1) -> List[str]:
    """
    Render first N pages of a PDF to PNG previews saved in upload_folder.
    Returns list of generated filenames (not full paths).
    """
    ensure_dir(upload_folder)
    filenames: List[str] = []
    try:
        doc = fitz.open(pdf_path)
        base = basename(pdf_path)
        total = min(pages, len(doc))
        for i in range(total):
            page = doc[i]
            # upscale for clarity
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            fname = f"{base}_p{i+1}_preview_{random_suffix()}.png"
            out_path = os.path.join(upload_folder, fname)
            pix.save(out_path)
            filenames.append(fname)
        doc.close()
    except Exception:
        return filenames
    return filenames


def render_pptx_previews(pptx_path: str, upload_folder: str, slides: int = 1) -> List[str]:
    """
    Best-effort preview for PPTX first N slides.
    Since python-pptx cannot render slides, we approximate by drawing text
    and pasting embedded images onto a canvas.
    Returns list of generated filenames (not full paths).
    """
    ensure_dir(upload_folder)
    filenames: List[str] = []
    try:
        prs = Presentation(pptx_path)
        base = basename(pptx_path)
        total = min(slides, len(prs.slides))
        # canvas size
        W, H = 1280, 720
        # default font
        try:
            font = ImageFont.truetype("arial.ttf", 24)
        except Exception:
            font = ImageFont.load_default()

        for s_idx in range(total):
            slide = prs.slides[s_idx]
            img = Image.new("RGB", (W, H), "white")
            draw = ImageDraw.Draw(img)
            y = 40

            # Draw text contents
            try:
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        text_lines: List[str] = []
                        for paragraph in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in paragraph.runs).strip()
                            if line:
                                text_lines.append(line)
                        if text_lines:
                            for line in text_lines[:8]:
                                draw.text((40, y), line, fill="black", font=font)
                                y += 32
                            y += 20
            except Exception:
                pass

            # Paste first picture if available
            try:
                for shape in slide.shapes:
                    if getattr(shape, "shape_type", None) == 13:  # picture
                        blob = shape.image.blob
                        from io import BytesIO
                        try:
                            pic = Image.open(BytesIO(blob)).convert("RGB")
                            # Fit picture into a box on the right
                            box_w, box_h = 560, 420
                            pic.thumbnail((box_w, box_h))
                            img.paste(pic, (W - box_w - 40, 40))
                            break
                        except Exception:
                            continue
            except Exception:
                pass

            fname = f"{base}_s{s_idx+1}_preview_{random_suffix()}.png"
            out_path = os.path.join(upload_folder, fname)
            img.save(out_path, format="PNG")
            filenames.append(fname)
    except Exception:
        return filenames

    return filenames